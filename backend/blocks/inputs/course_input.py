"""
blocks/inputs/course_input.py — Generic course material input block.

Auto-detects format (PDF, PPTX, video) and extracts text + page images.
Indexes text into Qdrant Cloud on execution.
"""
from __future__ import annotations

import uuid

from blocks.base import BlockData, BlockDescription, IBlock
from blocks.registry import register_block
from mixins.rag_mixin import RAGMixin


@register_block
class CourseInputBlock(IBlock, RAGMixin):
    description = BlockDescription(
        name="course_input",
        display_name="Course Material",
        group="input",
        input_types=[],  # source block — no upstream inputs
        output_types=["text"],
        parameters=[
            {
                "name": "language",
                "type": "select",
                "options": ["auto", "en", "fr", "ar"],
                "default": "auto",
                "label": "Transcription language (for video files)",
            },
            {
                "name": "extract_images",
                "type": "toggle",
                "default": True,
                "label": "Extract embedded images for infographic blocks",
            },
        ],
    )

    async def execute(self, inputs: list[BlockData], params: dict) -> list[BlockData]:
        """
        Inputs[0] must carry binary=<file bytes> and mime_type=<MIME>.
        If no inputs are provided (re-used from params), look for
        params["doc_id"] to reload from Qdrant.
        """
        if not inputs:
            return [BlockData(text="", metadata={"error": "No file input provided"})]

        bd = inputs[0]
        file_bytes = bd.binary
        mime = bd.mime_type or "application/octet-stream"
        user_id = bd.metadata.get("user_id", "anonymous")
        lang = params.get("language", "auto")
        extract_images = params.get("extract_images", True)

        if not file_bytes:
            return [BlockData(text="", metadata={"error": "No binary data in input"})]

        if mime == "application/pdf" or mime == "binary/pdf":
            text, images = self._extract_pdf(file_bytes, extract_images)
        elif "presentationml" in mime or mime.endswith(".pptx"):
            text, images = self._extract_pptx(file_bytes, extract_images)
        elif mime.startswith("video/"):
            text, images = await self._extract_video(file_bytes, lang, extract_images)
        else:
            # Treat as plain text
            text = file_bytes.decode("utf-8", errors="replace")
            images = []

        if not text.strip():
            return [BlockData(text="", metadata={"error": "No text could be extracted"})]

        # Chunk and index into Qdrant Cloud (best-effort)
        doc_id = bd.metadata.get("doc_id") or str(uuid.uuid4())
        from utils.pdf import chunk_text
        chunks = chunk_text(text, max_chars=1500)
        try:
            from db.qdrant import index_document
            index_document(doc_id, user_id, chunks)
        except Exception:
            pass  # Qdrant unavailable — continue without RAG

        slide_images_b64 = []
        if images:
            import base64
            slide_images_b64 = [base64.b64encode(img).decode() for img in images]

        return [BlockData(
            text=text,
            mime_type="text/plain",
            metadata={
                "doc_id": doc_id,
                "user_id": user_id,
                "page_count": len(images),
                "chunk_count": len(chunks),
                "slide_images": slide_images_b64[:20],  # first 20 pages max
            },
        )]

    # ── Format extractors ─────────────────────────────────────────────────────

    def _extract_pdf(
        self, data: bytes, extract_images: bool
    ) -> tuple[str, list[bytes]]:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=data, filetype="pdf")
        text = "\n\n---\n\n".join(
            page.get_text("text").strip() for page in doc if page.get_text().strip()
        )
        images: list[bytes] = []
        if extract_images:
            images = [page.get_pixmap(dpi=150).tobytes("png") for page in doc]
        doc.close()
        return text, images

    def _extract_pptx(
        self, data: bytes, extract_images: bool
    ) -> tuple[str, list[bytes]]:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return "[python-pptx not installed — PPTX extraction unavailable]", []

        import io
        prs = Presentation(io.BytesIO(data))
        slide_texts: list[str] = []
        for slide in prs.slides:
            parts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            slide_texts.append("\n".join(parts))
        text = "\n\n---\n\n".join(slide_texts)
        # Images: LibreOffice headless not guaranteed — skip for now
        return text, []

    async def _extract_video(
        self, data: bytes, lang: str, extract_images: bool
    ) -> tuple[str, list[bytes]]:
        import asyncio
        import os
        import tempfile

        try:
            from faster_whisper import WhisperModel  # type: ignore
        except ImportError:
            return "[faster-whisper not installed — video transcription unavailable]", []

        with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as f:
            f.write(data)
            tmp_path = f.name

        try:
            model = WhisperModel("tiny", device="cpu", compute_type="int8")
            effective_lang = None if lang == "auto" else lang
            segments, _ = await asyncio.get_event_loop().run_in_executor(
                None,
                lambda: model.transcribe(tmp_path, language=effective_lang),
            )
            text = " ".join(s.text for s in segments)
            images: list[bytes] = []
            if extract_images:
                images = self._extract_keyframes(tmp_path)
            return text, images
        finally:
            os.unlink(tmp_path)

    def _extract_keyframes(self, video_path: str) -> list[bytes]:
        """Extract one keyframe per 60 seconds using ffmpeg."""
        import subprocess
        import tempfile
        import os
        import glob

        with tempfile.TemporaryDirectory() as tmp:
            out_pattern = os.path.join(tmp, "frame_%04d.png")
            cmd = [
                "ffmpeg", "-i", video_path,
                "-vf", "fps=1/60",
                "-y", out_pattern,
            ]
            try:
                subprocess.run(cmd, capture_output=True, check=True)
            except Exception:
                return []
            frames = sorted(glob.glob(os.path.join(tmp, "frame_*.png")))
            images = []
            for fp in frames:
                with open(fp, "rb") as f:
                    images.append(f.read())
            return images
